# main.py

# --- 1. IMPORTS ---
import os
import time
import fitz  # PyMuPDF library for handling PDFs
import docx  # python-docx library for handling Word documents
import pptx  # python-pptx library for handling PowerPoint presentations
from PIL import Image
import pytesseract
from googletrans import Translator, LANGUAGES

# --- 2. CONFIGURATION ---
FOLDER_TO_WATCH = r'F:\image_translator\images_to_process'
PROCESSED_FOLDER = os.path.join(FOLDER_TO_WATCH, 'processed')
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# --- 3. HELPER FUNCTIONS ---

def setup_environment():
    """Checks if the required folders exist and creates them if they don't."""
    if not os.path.exists(FOLDER_TO_WATCH):
        print(f"Creating watch folder: '{FOLDER_TO_WATCH}'")
        os.makedirs(FOLDER_TO_WATCH)
    if not os.path.exists(PROCESSED_FOLDER):
        print(f"Creating processed folder: '{PROCESSED_FOLDER}'")
        os.makedirs(PROCESSED_FOLDER)
    print("-" * 30)
    print("Environment setup complete.")
    print(f"Watching for all supported file types in: {FOLDER_TO_WATCH}")
    print("Press Ctrl+C to stop the script.")
    print("-" * 30)

def extract_text_from_image(image_path_or_object):
    """Takes an image file path or an image object and uses Tesseract to extract text."""
    try:
        text = pytesseract.image_to_string(image_path_or_object)
        return text.strip()
    except pytesseract.TesseractNotFoundError:
        print(f"\nERROR: Tesseract executable not found. Please check the path: '{pytesseract.pytesseract.tesseract_cmd}'")
        return ""
    except Exception as e:
        print(f"Could not perform OCR. Error: {e}")
        return ""

def read_text_from_file(file_path):
    """Opens a plain text file (.txt) and reads its content."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read().strip()
    except Exception as e:
        print(f"Could not read text file '{os.path.basename(file_path)}'. Error: {e}")
        return None

def extract_text_from_pdf(file_path):
    """Extracts text from a PDF file, handling both text-based and image-based pages."""
    try:
        pdf_document = fitz.open(file_path)
        full_text = []
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            page_text = page.get_text().strip()
            if len(page_text) < 20:
                print(f"Page {page_num + 1} seems to be an image. Performing OCR...")
                pix = page.get_pixmap(dpi=300)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                ocr_text = extract_text_from_image(img)
                full_text.append(ocr_text)
            else:
                full_text.append(page_text)
        pdf_document.close()
        return "\n".join(full_text).strip()
    except Exception as e:
        print(f"Could not process PDF file '{os.path.basename(file_path)}'. Error: {e}")
        return None

def extract_text_from_docx(file_path):
    """Extracts text from a Microsoft Word (.docx) file."""
    try:
        document = docx.Document(file_path)
        full_text = [para.text for para in document.paragraphs]
        return "\n".join(full_text).strip()
    except Exception as e:
        print(f"Could not process Word file '{os.path.basename(file_path)}'. Error: {e}")
        return None

def extract_text_from_pptx(file_path):
    """Extracts text from a Microsoft PowerPoint (.pptx) file."""
    try:
        presentation = pptx.Presentation(file_path)
        full_text = []
        # Iterate through each slide in the presentation
        for slide in presentation.slides:
            # Iterate through each shape on the slide
            for shape in slide.shapes:
                # Check if the shape has a text frame
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text).strip()
    except Exception as e:
        print(f"Could not process PowerPoint file '{os.path.basename(file_path)}'. Error: {e}")
        return None

def translate_text_to_english(text):
    """Takes a string of text and translates it to English."""
    if not text:
        return None
    try:
        translator = Translator()
        translation_result = translator.translate(text, dest='en')
        source_language_name = LANGUAGES.get(translation_result.src, "Unknown")
        return {
            "translated_text": translation_result.text,
            "source_lang_name": source_language_name.capitalize()
        }
    except Exception as e:
        print(f"Translation failed. Could not connect to translation service. Error: {e}")
        return None

# --- 4. CORE PROCESSING FUNCTION ---
def process_files_in_folder():
    """Checks for files, processes one based on its type, and then moves it."""
    try:
        files = [f for f in os.listdir(FOLDER_TO_WATCH) if os.path.isfile(os.path.join(FOLDER_TO_WATCH, f))]
    except FileNotFoundError:
        print(f"ERROR: Watch folder not found at '{FOLDER_TO_WATCH}'. Exiting.")
        return

    if not files:
        return

    filename = files[0]
    file_path = os.path.join(FOLDER_TO_WATCH, filename)
    original_text = None

    print(f"\n>>> Found new file: '{filename}'")

    file_extension = os.path.splitext(filename)[1].lower()

    if file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']:
        print("File identified as an image. Performing OCR...")
        original_text = extract_text_from_image(file_path)
    elif file_extension == '.txt':
        print("File identified as a text file. Reading content...")
        original_text = read_text_from_file(file_path)
    elif file_extension == '.pdf':
        print("File identified as a PDF. Extracting text...")
        original_text = extract_text_from_pdf(file_path)
    elif file_extension == '.docx':
        print("File identified as a Word document. Extracting text...")
        original_text = extract_text_from_docx(file_path)
    elif file_extension == '.pptx':
        print("File identified as a PowerPoint presentation. Extracting text...")
        original_text = extract_text_from_pptx(file_path)
    else:
        print(f"Unsupported file type: '{file_extension}'. This file will be moved without processing.")

    if original_text:
        print(f"\n--- Original Text ---\n{original_text}\n---------------------")
        translation_info = translate_text_to_english(original_text)
        if translation_info:
            print(f"Detected Language: {translation_info['source_lang_name']}")
            print(f"\n--- Translated Text (English) ---\n{translation_info['translated_text']}\n---------------------------------")
    else:
        print("No text was extracted from the file.")

    try:
        destination_path = os.path.join(PROCESSED_FOLDER, filename)
        os.rename(file_path, destination_path)
        print(f"Moved '{filename}' to the 'processed' folder.")
    except Exception as e:
        print(f"ERROR: Could not move file '{filename}'. It may be in use. Error: {e}")

# --- 5. MAIN EXECUTION BLOCK ---
if __name__ == "__main__":
    setup_environment()
    try:
        while True:
            process_files_in_folder()
            time.sleep(5)
    except KeyboardInterrupt:
        print("\nScript stopped by user. Goodbye!")